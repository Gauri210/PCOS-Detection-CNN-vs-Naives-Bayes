import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
import os

prs = Presentation()

# Helper function to add title and content slides
def add_slide(prs, title, content_points, image_path=None, image_layout='right'):
    # Choose slide layout: 1 is title and content
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Set content
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    
    for i, point in enumerate(content_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.space_after = Pt(14)
        
    if image_path and os.path.exists(image_path):
        if image_layout == 'right':
            # Resize text box
            body_shape.width = Inches(4.5)
            # Add image
            slide.shapes.add_picture(image_path, Inches(5), Inches(2), width=Inches(4.5))
        elif image_layout == 'bottom':
            # Adjust text box
            body_shape.height = Inches(2.5)
            # Add image at bottom
            slide.shapes.add_picture(image_path, Inches(1), Inches(4), width=Inches(8))
            
    return slide

# Title Slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Automated PCOS Detection from Ultrasound Imagery"
subtitle.text = "A Comparative Analysis of Naive Bayes and Convolutional Neural Networks\n\nAuthor: Gauri\nApril 2026"

# Slide 1: Introduction
add_slide(prs, "I. Introduction", [
    "Polycystic Ovary Syndrome (PCOS) is a common hormonal disorder among women.",
    "Early detection through ultrasound imagery is crucial for effective treatment.",
    "Manual interpretation by doctors can be subjective and time-consuming.",
    "Goal: Build an automated computer program to classify ovarian ultrasound images as Normal or Abnormal."
])

# Slide 2: Methodology - Dataset Configuration
add_slide(prs, "II. Dataset Information", [
    "Custom dataset split into training (3,200 images) and test sets (1,468 images).",
    "Class 0: Normal ovaries (903 training images).",
    "Class 1: Abnormal / PCOS ovaries (2,297 training images).",
    "The dataset reflects real-world clinical imbalances."
], "graph_cell_6_0.png", "right")

# Slide 2.1: Class Distribution
add_slide(prs, "III. Class Distribution Profile", [
    "Significant categorical imbalance inherent in the clinical repository.",
    "Normal cases: 903 | Abnormal cases: 2,297.",
    "Requires robust metrics (like Recall) to evaluate effectively."
], "graph_cell_4_0.png", "right")

# Slide 3: Tensor Standardization and Spatial Preprocessing
add_slide(prs, "IV. Image Preprocessing", [
    "1. Grayscale Conversion: Removes unnecessary color information.",
    "2. Resizing: Resizes all images to 128x128 pixels for uniformity.",
    "3. Normalization: Scales pixel values to a standard range (Z-score normalization).",
    "4. Flattening: Converts 2D images into 1D vectors for the Naive Bayes model."
])

# Slide 3.1: Structural Difference Heatmap
add_slide(prs, "V. Structural Difference Heatmap", [
    "Visual structural comparison of Average Normal vs Average Abnormal Ultrasounds.",
    "Heatmap exposes subtle regions of morphological variance.",
    "Confirms that pathogenic features (follicular distributions) are spatially correlated."
], "graph_cell_10_2.png", "bottom")

# Slide 3.2: Pixel Intensity Distribution
add_slide(prs, "VI. Pixel Intensity Distributions", [
    "Kernel Density Estimation across Normal and Abnormal classes.",
    "Significant overlap demonstrates non-linear structural complexities.",
    "Highlights why baseline probabilistic classification (Naive Bayes) will struggle."
], "graph_cell_12_3.png", "right")

# Slide 4: Baseline Probabilistic Formulation
add_slide(prs, "VII. Gaussian Naive Bayes Model", [
    "Serves as the traditional machine learning baseline.",
    "Treats every pixel independently without considering its neighbors.",
    "Limited because it ignores the spatial shapes and patterns in the image.",
    "Struggles to reliably identify clustered structures like PCOS follicles."
])

# Slide 5: Convolutional Backbone Architecture
add_slide(prs, "VIII. Convolutional Neural Network (CNN)", [
    "Designed to exploit local spatial correlations within a grid-like topology.",
    "Convolutional Blocks: 4 layers with progressively increasing filters (32 -> 256) utilizing 3x3 kernels.",
    "Activation & Pooling: ReLU activations coupled with 2x2 Max Pooling.",
    "Regularization: Batch Normalization and 50% Dropout regularization.",
    "Output: Softmax dense layer for final binary classification."
])

# Slide 6: Empirical Evaluation - Categorical Assessment
add_slide(prs, "IX. Results and Evaluation", [
    "Focus lies heavily on Recall (Sensitivity) given the clinical context of medical screening.",
    "Minimizing False Negatives ensures critical pathological structures are not overlooked.",
    "Performance validated across standard categorical dimensions: Accuracy, Precision, Recall, and F1-Score."
])

# Slide 7: Gaussian Naive Bayes Performance
add_slide(prs, "X. Gaussian Naive Bayes Performance", [
    "Accuracy: 61.19%",
    "Precision: 86.20%",
    "Recall: 54.68%",
    "F1-Score: 66.92%",
    "Critically hampered by False Negatives due to feature independence assumptions."
], "graph_cell_14_1.png", "right")

# Slide 8: CNN Performance Metrics
add_slide(prs, "XI. Convolutional Network Performance", [
    "Accuracy: 78.78%",
    "Precision: 77.59%",
    "Recall: 99.04%",
    "F1-Score: 87.01%",
    "Demonstrates robust spatial feature extraction and minimal False Negatives."
], "graph_cell_26_3.png", "right")

# Slide 9: Optimization Stability
add_slide(prs, "XII. CNN Optimization Stability", [
    "Trained over 50 epochs utilizing the Adam optimizer.",
    "Steady convergence patterns across loss gradients without aggressive geometric spiking.",
    "Excellent generalization verified by validation metrics."
], "graph_cell_23_2.png", "bottom")

# Slide 10: Comparative Diagnostic Analysis
add_slide(prs, "XIII. Performance Comparison", [
    "The CNN yielded a massive +17.59% absolute increase in Accuracy.",
    "Crucially, the CNN delivered a +44.36% surge in Recall capability.",
    "The deep spatial feature mapping drastically outperformed probabilistic, pixel-wise estimation."
], "graph_cell_31_4.png", "right")

# Slide 11: Concluding Analysis and Future Developments
add_slide(prs, "XIV. Conclusion and Future Work", [
    "Deep convolutional architectures are fundamentally better equipped for CAD systems in gynecology.",
    "Successfully bypassed the theoretical limitations of pixel-wise probability mapping.",
    "Future vectors: Integrating Transfer Learning (EfficientNet, Vision Transformers).",
    "Future vectors: Deploying spatial attention mechanisms (like Grad-CAM) to establish visible structural accountability."
])

prs.save("PCOS_Detection_Presentation.pptx")
